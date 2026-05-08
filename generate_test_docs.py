from docx import Document
from pptx import Presentation

# Generate test.docx
doc = Document()
doc.add_heading('Front Matter', 0)
doc.add_paragraph('This is a test document.')

doc.add_heading('Abstract', 1)
doc.add_paragraph('This is the abstract for the test document.')

doc.add_heading('1 Introduction', 1)
doc.add_paragraph('This is the introduction section. It has some text.')

doc.add_heading('2 Methodology', 1)
doc.add_paragraph('We used python-docx to generate this.')

doc.save('test.docx')

# Generate test.pptx
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Front Matter"
subtitle.text = "This is a test presentation."

bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Abstract'
tf = body_shape.text_frame
tf.text = 'This is the abstract for the test presentation.'

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = '1 Introduction'
tf = body_shape.text_frame
tf.text = 'This is the introduction section. It has some text.'

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = '2 Methodology'
tf = body_shape.text_frame
tf.text = 'We used python-pptx to generate this.'

prs.save('test.pptx')
