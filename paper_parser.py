import re
from typing import List, Dict
from pypdf import PdfReader

SECTION_HEADING_KEYWORDS = [
    'Abstract', 'Introduction', 'Background', 'Related Work', 'Methodology',
    'Methods', 'Experiments', 'Results', 'Discussion', 'Conclusion',
    'References', 'Appendix', 'Training', 'Attention', 'Optimizer',
    'Regularization', 'Hardware', 'Evaluation'
]

HEADING_PATTERN = re.compile(
    r'^(?:\d+(?:\.\d+)*[\.\s]*)?(?:' +
    '|'.join(re.escape(keyword) for keyword in SECTION_HEADING_KEYWORDS) +
    r')\s*$',
    re.IGNORECASE
)
NUMBERED_HEADING_PATTERN = re.compile(r'^\d+(?:\.\d+)*[\.\s]+[A-Z][A-Za-z0-9&\-\s]{0,100}$')
SECTION_NUMBER_PATTERN = re.compile(r'^\d+(?:\.\d+)*$')
SHORT_TITLE_PATTERN = re.compile(r'^[A-Za-z][A-Za-z0-9&\-\s]{0,60}$')


def is_heading_line(line: str) -> bool:
    if HEADING_PATTERN.match(line):
        return True
    if NUMBERED_HEADING_PATTERN.match(line):
        return True
    return False


def should_combine_number_and_title(line: str, next_line: str) -> bool:
    if SECTION_NUMBER_PATTERN.match(line):
        return bool(next_line and SHORT_TITLE_PATTERN.match(next_line) and len(next_line.split()) <= 7)
    return False


def extract_lines_from_pdf(pdf_path: str) -> List[str]:
    reader = PdfReader(pdf_path)
    lines = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            for line in text.split('\n'):
                if line.strip():
                    lines.append(line.strip())
    return lines

def extract_lines_from_docx(docx_path: str) -> List[str]:
    from docx import Document
    doc = Document(docx_path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            for line in text.split('\n'):
                if line.strip():
                    lines.append(line.strip())
    return lines

def extract_lines_from_pptx(pptx_path: str) -> List[str]:
    from pptx import Presentation
    prs = Presentation(pptx_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    for line in text.split('\n'):
                        if line.strip():
                            lines.append(line.strip())
    return lines

def parse_sections_from_lines(raw_lines: List[str]) -> List[Dict[str, str]]:
    sections: List[Dict[str, str]] = []
    current_heading = "Front Matter"
    current_content: List[str] = []

    i = 0
    while i < len(raw_lines):
        line = raw_lines[i]

        if i + 1 < len(raw_lines) and should_combine_number_and_title(line, raw_lines[i + 1]):
            line = f"{line} {raw_lines[i + 1]}"
            i += 1

        if is_heading_line(line):
            if current_content:
                sections.append({
                    'heading': current_heading,
                    'content': ' '.join(current_content).strip()
                })
            current_heading = line
            current_content = []
        else:
            if not SECTION_NUMBER_PATTERN.match(line):
                current_content.append(line)

        i += 1

    if current_content:
        sections.append({
            'heading': current_heading,
            'content': ' '.join(current_content).strip()
        })

    return sections

def extract_sections_from_document(file_path: str) -> List[Dict[str, str]]:
    ext = file_path.lower().split('.')[-1]
    if ext == 'pdf':
        lines = extract_lines_from_pdf(file_path)
    elif ext == 'docx':
        lines = extract_lines_from_docx(file_path)
    elif ext == 'pptx':
        lines = extract_lines_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")
    
    return parse_sections_from_lines(lines)


if __name__ == "__main__":
    import sys
    sections = extract_sections_from_document(sys.argv[1])
    for s in sections[:10]:
        print(f"=== {s['heading']} ===\n{s['content'][:200]}...\n")